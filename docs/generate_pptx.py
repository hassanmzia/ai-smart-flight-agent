#!/usr/bin/env python3
"""
Generate a professional multi-slide Technical Architecture PowerPoint
for the AI Smart Flight Agent project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────
DARK_BG      = RGBColor(0x1B, 0x1F, 0x3B)   # deep navy
MED_BG       = RGBColor(0x27, 0x2D, 0x52)   # medium navy
ACCENT_BLUE  = RGBColor(0x00, 0x7A, 0xCC)   # bright blue
ACCENT_CYAN  = RGBColor(0x00, 0xBC, 0xD4)   # cyan
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)   # green
ACCENT_ORANGE= RGBColor(0xFF, 0x98, 0x00)   # orange
ACCENT_RED   = RGBColor(0xE5, 0x39, 0x35)   # red
ACCENT_PURPLE= RGBColor(0x7E, 0x57, 0xC2)   # purple
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT     = RGBColor(0x1A, 0x1A, 0x2E)
SOFT_WHITE   = RGBColor(0xF0, 0xF0, 0xF5)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

# ── Helpers ─────────────────────────────────────────────────────────

def _solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    ln = shape.line
    if border_color:
        ln.color.rgb = border_color
        ln.width = border_width
    else:
        ln.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _set_paragraph(tf, text, font_size=12, bold=False, color=WHITE,
                   alignment=PP_ALIGN.LEFT, font_name="Calibri", space_after=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    return p


def _add_arrow(slide, x1, y1, x2, y2, color=ACCENT_CYAN, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1,  # straight connector
        x1, y1, x2, y2
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # Add end arrow
    connector.end_x = x2
    connector.end_y = y2
    return connector


def _slide_title_bar(slide, title_text):
    """Dark header bar at the top of content slides."""
    bar = _add_rect(slide, Inches(0), Inches(0), SW, Inches(1.0), ACCENT_BLUE)
    _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                 title_text, font_size=28, bold=True, color=WHITE, font_name="Calibri")
    return bar


def _box_with_header(slide, left, top, width, height, header, body_lines,
                     header_color=ACCENT_BLUE, body_bg=MED_BG, text_color=WHITE,
                     header_font=12, body_font=10):
    """Rounded rect with a coloured header strip and bullet body."""
    # body
    body_shape = _add_rect(slide, left, top, width, height, body_bg,
                           border_color=header_color, border_width=Pt(1.5))
    # header strip
    _add_rect(slide, left, top, width, Inches(0.42), header_color)
    _add_textbox(slide, left + Inches(0.1), top + Inches(0.04),
                 width - Inches(0.2), Inches(0.38),
                 header, font_size=header_font, bold=True, color=WHITE)
    # body text
    if body_lines:
        txBox = slide.shapes.add_textbox(left + Inches(0.12),
                                          top + Inches(0.48),
                                          width - Inches(0.24),
                                          height - Inches(0.55))
        tf = txBox.text_frame
        tf.word_wrap = True
        first = True
        for line in body_lines:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(body_font)
            p.font.color.rgb = text_color
            p.font.name = "Calibri"
            p.space_after = Pt(2)
    return body_shape


# ====================================================================
# SLIDE 1 – Title
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_solid_bg(slide, DARK_BG)

# decorative accent bar
_add_rect(slide, Inches(0), Inches(0), SW, Inches(0.15), ACCENT_BLUE)
_add_rect(slide, Inches(0), Inches(7.35), SW, Inches(0.15), ACCENT_BLUE)

# Left accent stripe
_add_rect(slide, Inches(0.8), Inches(1.8), Inches(0.08), Inches(2.2), ACCENT_CYAN)

# Title
_add_textbox(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(1.2),
             "AI Smart Flight Agent", font_size=48, bold=True, color=WHITE)

# Subtitle
_add_textbox(slide, Inches(1.2), Inches(3.0), Inches(10), Inches(0.7),
             "Technical Architecture Overview", font_size=28, bold=False, color=ACCENT_CYAN)

# Separator line
_add_rect(slide, Inches(1.2), Inches(3.8), Inches(4), Inches(0.04), ACCENT_BLUE)

# Footer
_add_textbox(slide, Inches(1.2), Inches(4.2), Inches(10), Inches(0.5),
             "Multi-Agent AI Travel Planning System  |  Docker Compose Deployment",
             font_size=16, bold=False, color=LIGHT_GRAY)

# decorative boxes on right
for i, clr in enumerate([ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE]):
    _add_rect(slide, Inches(9.5 + i * 0.9), Inches(5.2), Inches(0.7), Inches(0.7), clr)


# ====================================================================
# SLIDE 2 – System Architecture Overview
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, DARK_BG)
_slide_title_bar(slide, "System Architecture Overview")

layer_data = [
    ("Client Layer",       ["Browser", "Mobile", "API Clients"],            ACCENT_PURPLE),
    ("Reverse Proxy",      ["Nginx", "Port 80 / 443", "SSL Termination"],  ACCENT_BLUE),
    ("Application Layer",  ["Frontend :3090", "Backend :8109", "MCP :8107"], ACCENT_CYAN),
    ("Data Layer",         ["PostgreSQL :5438", "Redis :6384", "RabbitMQ :5673"], ACCENT_GREEN),
    ("External Services",  ["OpenAI API", "SerpAPI", "Stripe / Weather"],  ACCENT_ORANGE),
]

box_w = Inches(2.2)
box_h = Inches(1.6)
gap   = Inches(0.22)
total = len(layer_data) * box_w + (len(layer_data) - 1) * gap
start_x = (SW - total) // 2
y_pos   = Inches(2.4)

for idx, (header, lines, clr) in enumerate(layer_data):
    x = start_x + idx * (box_w + gap)
    _box_with_header(slide, x, y_pos, box_w, box_h, header, lines,
                     header_color=clr, body_font=11)

# Arrows between layers
arrow_y = y_pos + box_h // 2
for idx in range(len(layer_data) - 1):
    x1 = start_x + (idx + 1) * (box_w + gap) - gap + Inches(0.02)
    x2 = x1 + gap - Inches(0.04)
    slide.shapes.add_connector(1, x1, arrow_y, x2, arrow_y).line.color.rgb = ACCENT_CYAN
    # arrowhead triangle
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x2 - Inches(0.08),
                                  arrow_y - Inches(0.08), Inches(0.16), Inches(0.16))
    tri.rotation = 90.0
    tri.fill.solid()
    tri.fill.fore_color.rgb = ACCENT_CYAN
    tri.line.fill.background()

# Flow label
_add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(0.4),
             "Request Flow:  Client  -->  Nginx  -->  App Servers  -->  Data Stores  -->  External APIs",
             font_size=14, bold=False, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom note boxes
note_data = [
    ("Docker Compose", "8 containers on\ntravel-network bridge", ACCENT_BLUE),
    ("Health Checks", "All services monitored\nrestart: unless-stopped", ACCENT_GREEN),
    ("Volumes", "postgres_data, redis_data\nrabbitmq_data, static, media", ACCENT_ORANGE),
]
for i, (hdr, body, clr) in enumerate(note_data):
    _box_with_header(slide, Inches(1.2) + i * Inches(3.9), Inches(5.2),
                     Inches(3.5), Inches(1.1), hdr, body.split("\n"),
                     header_color=clr, body_font=10)


# ====================================================================
# SLIDE 3 – Frontend Architecture
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, DARK_BG)
_slide_title_bar(slide, "Frontend Architecture")

# Core stack box
_box_with_header(slide, Inches(0.4), Inches(1.3), Inches(3.8), Inches(2.6),
                 "Core Stack", [
                     "React 18 + TypeScript + Vite",
                     "State: Zustand (Auth, Booking,",
                     "  Notification, Search stores)",
                     "Server State: React Query + Axios",
                     "  with JWT interceptor",
                     "Styling: TailwindCSS + Dark Mode",
                     "Mobile-responsive design",
                 ], header_color=ACCENT_BLUE, body_font=10)

# Pages box
_box_with_header(slide, Inches(4.5), Inches(1.3), Inches(4.3), Inches(2.6),
                 "32 Pages", [
                     "Home, AI Planner, Flights, Hotels",
                     "Cars, Restaurants, Weather, Events",
                     "Shopping, Safety, Dashboard, Bookings",
                     "Itinerary, Payment, Reviews, Analytics",
                     "Profile, Settings, Login, Register",
                     "Tourist Attractions, Commute Planner",
                     "Admin, Error Pages, and more...",
                 ], header_color=ACCENT_CYAN, body_font=10)

# Services & Widget box
_box_with_header(slide, Inches(9.1), Inches(1.3), Inches(3.8), Inches(2.6),
                 "Services & Components", [
                     "20 Service Modules (API clients)",
                     "Floating AI Chat Widget",
                     "  - Voice input support",
                     "  - Real-time streaming",
                     "WebSocket notifications",
                     "JWT auto-refresh interceptor",
                     "Responsive breakpoints",
                 ], header_color=ACCENT_GREEN, body_font=10)

# Architecture diagram row
stores = ["authStore", "bookingStore", "notificationStore", "searchStore"]
sx = Inches(0.4)
for i, s in enumerate(stores):
    clr = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE][i]
    _box_with_header(slide, sx + i * Inches(2.3), Inches(4.3), Inches(2.1), Inches(0.85),
                     "Zustand", [s], header_color=clr, body_font=10)

# React Query layer
_add_rect(slide, Inches(9.7), Inches(4.3), Inches(3.2), Inches(0.85), MED_BG,
          border_color=ACCENT_PURPLE, border_width=Pt(1.5))
_add_rect(slide, Inches(9.7), Inches(4.3), Inches(3.2), Inches(0.42), ACCENT_PURPLE)
_add_textbox(slide, Inches(9.8), Inches(4.34), Inches(3.0), Inches(0.38),
             "React Query", font_size=12, bold=True, color=WHITE)
_add_textbox(slide, Inches(9.8), Inches(4.78), Inches(3.0), Inches(0.3),
             "Cache + Axios + JWT", font_size=10, color=WHITE)

# Tech badges row
tech_badges = [
    ("React 18", ACCENT_BLUE), ("TypeScript 5", ACCENT_CYAN),
    ("Vite", ACCENT_GREEN), ("TailwindCSS 3", ACCENT_PURPLE),
    ("Zustand", ACCENT_ORANGE), ("React Query", ACCENT_RED),
    ("Axios", ACCENT_BLUE), ("WebSocket", ACCENT_CYAN),
]
for i, (label, clr) in enumerate(tech_badges):
    bx = Inches(0.4) + i * Inches(1.6)
    shape = _add_rect(slide, bx, Inches(5.5), Inches(1.45), Inches(0.42), clr)
    _add_textbox(slide, bx, Inches(5.53), Inches(1.45), Inches(0.38),
                 label, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Port info
_add_textbox(slide, Inches(0.4), Inches(6.2), Inches(12), Inches(0.4),
             "Dev Server: Vite on port 3090  |  Build: Optimized static assets served by Nginx  |  Env: VITE_API_BASE_URL, VITE_WS_URL",
             font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)


# ====================================================================
# SLIDE 4 – Backend Architecture
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, DARK_BG)
_slide_title_bar(slide, "Backend Architecture")

# Core stack
_box_with_header(slide, Inches(0.4), Inches(1.3), Inches(4.0), Inches(2.5),
                 "Core Stack", [
                     "Django 5.0 + DRF",
                     "Gunicorn (4 workers)",
                     "Authentication: JWT (SimpleJWT)",
                     "  Access: 1h  |  Refresh: 7d",
                     "WebSocket: Django Channels + Daphne",
                     "API Docs: drf-spectacular",
                     "  (OpenAPI / Swagger UI)",
                 ], header_color=ACCENT_BLUE, body_font=10)

# Django Apps
apps = [
    "users", "agents", "flights", "hotels", "bookings", "payments",
    "itineraries", "notifications", "reviews", "analytics",
    "car_rentals", "restaurants", "attractions", "weather",
    "events", "shopping", "safety", "commute", "tourist_attractions"
]
_box_with_header(slide, Inches(4.7), Inches(1.3), Inches(4.3), Inches(2.5),
                 "19 Django Apps", [
                     ", ".join(apps[:5]),
                     ", ".join(apps[5:10]),
                     ", ".join(apps[10:14]),
                     ", ".join(apps[14:]),
                 ], header_color=ACCENT_CYAN, body_font=10)

# Background tasks
_box_with_header(slide, Inches(9.3), Inches(1.3), Inches(3.6), Inches(2.5),
                 "Background Tasks", [
                     "Celery Worker",
                     "  4 concurrent workers",
                     "Celery Beat",
                     "  DB-backed Scheduler",
                     "Broker: RabbitMQ (AMQP)",
                     "Result Backend: Redis",
                     "Task retry with backoff",
                 ], header_color=ACCENT_ORANGE, body_font=10)

# Middleware / Security row
_box_with_header(slide, Inches(0.4), Inches(4.1), Inches(6.2), Inches(1.4),
                 "Middleware & Security", [
                     "CORS (allowed origins)  |  CSRF Protection  |  Secure Cookies",
                     "SSL Redirect  |  HSTS 1yr  |  XSS Filter  |  X-Frame DENY",
                     "Content-Type Nosniff  |  Session Security  |  Rate Limiting",
                 ], header_color=ACCENT_RED, body_font=10)

# API structure
_box_with_header(slide, Inches(6.9), Inches(4.1), Inches(5.9), Inches(1.4),
                 "API Endpoints Structure", [
                     "/api/auth/   |  /api/users/   |  /api/flights/  |  /api/hotels/",
                     "/api/agents/  |  /api/bookings/  |  /api/payments/  |  /api/itineraries/",
                     "/api/reviews/  |  /api/analytics/  |  /api/notifications/  |  /api/weather/",
                 ], header_color=ACCENT_GREEN, body_font=10)

# Tech badges
be_badges = [
    ("Django 5.0", ACCENT_BLUE), ("DRF", ACCENT_CYAN), ("Gunicorn", ACCENT_GREEN),
    ("Celery", ACCENT_ORANGE), ("Channels", ACCENT_PURPLE), ("SimpleJWT", ACCENT_RED),
    ("PostgreSQL", ACCENT_BLUE), ("Redis", ACCENT_RED),
]
for i, (label, clr) in enumerate(be_badges):
    bx = Inches(0.4) + i * Inches(1.6)
    _add_rect(slide, bx, Inches(5.8), Inches(1.45), Inches(0.42), clr)
    _add_textbox(slide, bx, Inches(5.83), Inches(1.45), Inches(0.38),
                 label, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

_add_textbox(slide, Inches(0.4), Inches(6.5), Inches(12), Inches(0.4),
             "Server: Port 8109  |  Python 3.11  |  Logging: RotatingFileHandler (10 MB x 10)  |  Static: WhiteNoise + Nginx",
             font_size=11, color=LIGHT_GRAY)


# ====================================================================
# SLIDE 5 – Multi-Agent AI System
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, DARK_BG)
_slide_title_bar(slide, "Multi-Agent AI System (LangGraph)")

# LLM Config
_box_with_header(slide, Inches(0.4), Inches(1.3), Inches(3.0), Inches(1.5),
                 "LLM Configuration", [
                     "Model: OpenAI GPT-4o-mini",
                     "Temperature: 0.7",
                     "Provider: ChatOpenAI",
                     "Streaming: Enabled",
                 ], header_color=ACCENT_PURPLE, body_font=10)

# Shared state
_box_with_header(slide, Inches(3.7), Inches(1.3), Inches(3.3), Inches(1.5),
                 "Shared State", [
                     "TravelAgentState",
                     "  - messages[]",
                     "  - search_results{}",
                     "  - evaluations{}",
                 ], header_color=ACCENT_BLUE, body_font=10)

# Enhanced orchestrator
_box_with_header(slide, Inches(7.3), Inches(1.3), Inches(5.6), Inches(1.5),
                 "Enhanced Orchestrator", [
                     "RAG Pipeline  |  ThreadPoolExecutor(8 workers)",
                     "Sub-agents: Health, Visa, Packing, LocalExpert",
                     "Redis Caching  |  Parallel execution  |  Error recovery",
                 ], header_color=ACCENT_ORANGE, body_font=10)

# Sequential pipeline
pipeline_label = "Sequential Pipeline (StateGraph)"
_add_textbox(slide, Inches(0.4), Inches(3.05), Inches(6), Inches(0.4),
             pipeline_label, font_size=16, bold=True, color=ACCENT_CYAN)

# Search agents row
search_agents = [
    ("Flight\nAgent", ACCENT_BLUE),
    ("Hotel\nAgent", ACCENT_CYAN),
    ("Car Rental\nAgent", ACCENT_GREEN),
    ("Restaurant\nAgent", ACCENT_ORANGE),
]
y_search = Inches(3.55)
for i, (name, clr) in enumerate(search_agents):
    x = Inches(0.4) + i * Inches(2.1)
    shape = _add_rect(slide, x, y_search, Inches(1.85), Inches(0.85), clr)
    _add_textbox(slide, x, y_search + Inches(0.1), Inches(1.85), Inches(0.7),
                 name, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Arrow triangles between search agents
for i in range(3):
    ax = Inches(0.4) + (i + 1) * Inches(2.1) - Inches(0.2)
    tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  ax, y_search + Inches(0.28), Inches(0.35), Inches(0.3))
    tri.fill.solid()
    tri.fill.fore_color.rgb = WHITE
    tri.line.fill.background()

# Eval agents row
_add_textbox(slide, Inches(0.4), Inches(4.6), Inches(6), Inches(0.35),
             "Evaluation Agents:", font_size=13, bold=True, color=ACCENT_CYAN)

eval_agents = [
    ("Goal-Based\nEvaluator", ACCENT_PURPLE),
    ("Utility-Based\nEvaluator", ACCENT_RED),
    ("Car\nEvaluator", ACCENT_GREEN),
    ("Restaurant\nEvaluator", ACCENT_ORANGE),
]
y_eval = Inches(5.0)
for i, (name, clr) in enumerate(eval_agents):
    x = Inches(0.4) + i * Inches(2.1)
    shape = _add_rect(slide, x, y_eval, Inches(1.85), Inches(0.85), clr)
    _add_textbox(slide, x, y_eval + Inches(0.1), Inches(1.85), Inches(0.7),
                 name, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

for i in range(3):
    ax = Inches(0.4) + (i + 1) * Inches(2.1) - Inches(0.2)
    tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  ax, y_eval + Inches(0.28), Inches(0.35), Inches(0.3))
    tri.fill.solid()
    tri.fill.fore_color.rgb = WHITE
    tri.line.fill.background()

# Manager → END
_add_rect(slide, Inches(8.8), Inches(4.2), Inches(1.8), Inches(0.9), ACCENT_RED)
_add_textbox(slide, Inches(8.8), Inches(4.35), Inches(1.8), Inches(0.7),
             "Manager\nAgent", font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                              Inches(10.7), Inches(4.45), Inches(0.4), Inches(0.35))
tri.fill.solid()
tri.fill.fore_color.rgb = WHITE
tri.line.fill.background()

_add_rect(slide, Inches(11.2), Inches(4.2), Inches(1.4), Inches(0.9), ACCENT_GREEN)
_add_textbox(slide, Inches(11.2), Inches(4.35), Inches(1.4), Inches(0.7),
             "END", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Tools box
_box_with_header(slide, Inches(8.8), Inches(5.4), Inches(3.8), Inches(1.2),
                 "Agent Tools", [
                     "SerpAPI: Google Flights",
                     "SerpAPI: Google Hotels",
                     "SerpAPI: Google Local (Cars/Food)",
                 ], header_color=ACCENT_BLUE, body_font=10)

_add_textbox(slide, Inches(0.4), Inches(6.2), Inches(12), Inches(0.4),
             "Pipeline: Search Phase (4 agents)  -->  Evaluation Phase (4 agents)  -->  Manager Agent  -->  Final Recommendation",
             font_size=11, color=LIGHT_GRAY)


# ====================================================================
# SLIDE 6 – MCP Server
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, DARK_BG)
_slide_title_bar(slide, "MCP Agent Communication Server")

# Server core
_box_with_header(slide, Inches(0.4), Inches(1.3), Inches(3.8), Inches(2.2),
                 "Server Core", [
                     "FastAPI + Uvicorn",
                     "Port: 8107",
                     "Async request handling",
                     "Health check endpoint",
                     "CORS enabled",
                 ], header_color=ACCENT_BLUE, body_font=11)

# Agent Registration
_box_with_header(slide, Inches(4.5), Inches(1.3), Inches(4.0), Inches(2.2),
                 "Agent Registration", [
                     "Flight Agent",
                     "Hotel Agent",
                     "Goal-Based Evaluator",
                     "Utility-Based Evaluator",
                     "Manager Agent",
                 ], header_color=ACCENT_CYAN, body_font=11)

# Message types
_box_with_header(slide, Inches(8.8), Inches(1.3), Inches(4.1), Inches(2.2),
                 "Message Types", [
                     "Request  (agent-to-agent)",
                     "Response (result payload)",
                     "Notification (broadcast)",
                     "Error    (fault handling)",
                 ], header_color=ACCENT_ORANGE, body_font=11)

# Communication layer
_box_with_header(slide, Inches(0.4), Inches(3.9), Inches(6.0), Inches(2.0),
                 "Communication Layer", [
                     "WebSocket endpoint: /ws/{agent_id}",
                     "Redis Pub/Sub queues for inter-agent messaging",
                     "Async message dispatching",
                     "Connection pooling & reconnection logic",
                     "Message serialization (JSON)",
                 ], header_color=ACCENT_PURPLE, body_font=11)

# Shared context
_box_with_header(slide, Inches(6.7), Inches(3.9), Inches(5.9), Inches(2.0),
                 "Shared Context & State", [
                     "Session-scoped TTL for context data",
                     "Redis-backed state persistence",
                     "Agent capability discovery",
                     "Load balancing across agent instances",
                     "Graceful shutdown & cleanup",
                 ], header_color=ACCENT_GREEN, body_font=11)

_add_textbox(slide, Inches(0.4), Inches(6.3), Inches(12), Inches(0.4),
             "Flow: Agent registers --> Receives agent_id --> Connects via WebSocket --> Sends/receives messages via Redis Pub/Sub",
             font_size=11, color=LIGHT_GRAY)


# ====================================================================
# SLIDE 7 – Data Layer
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, DARK_BG)
_slide_title_bar(slide, "Data & Messaging Layer")

# PostgreSQL
_box_with_header(slide, Inches(0.4), Inches(1.3), Inches(4.0), Inches(2.5),
                 "PostgreSQL 15  (Port 5438)", [
                     "Users & Authentication",
                     "Bookings & Reservations",
                     "Payments & Transactions",
                     "Itineraries & Plans",
                     "Reviews & Ratings",
                     "Analytics & Metrics",
                     "All Django model data",
                 ], header_color=ACCENT_BLUE, body_font=11)

# Redis
_box_with_header(slide, Inches(4.7), Inches(1.3), Inches(4.0), Inches(2.5),
                 "Redis 7  (Port 6384)", [
                     "DB 0: Django Cache",
                     "  - View/query caching",
                     "  - Session storage",
                     "DB 1: Channels + MCP",
                     "  - WebSocket channel layers",
                     "  - MCP agent message queues",
                     "  - Pub/Sub for real-time",
                 ], header_color=ACCENT_RED, body_font=11)

# RabbitMQ
_box_with_header(slide, Inches(9.0), Inches(1.3), Inches(3.9), Inches(2.5),
                 "RabbitMQ 3.12", [
                     "AMQP Port: 5673",
                     "Management UI: 15673",
                     "Celery task broker",
                     "Task queues & routing",
                     "Dead letter exchanges",
                     "Message persistence",
                     "Monitoring dashboard",
                 ], header_color=ACCENT_ORANGE, body_font=11)

# Docker volumes
_box_with_header(slide, Inches(0.4), Inches(4.2), Inches(6.0), Inches(1.8),
                 "Docker Volumes", [
                     "postgres_data   - Database files",
                     "redis_data      - Redis persistence (AOF/RDB)",
                     "rabbitmq_data   - Queue data & config",
                     "static_volume   - Django static files",
                     "media_volume    - User uploads",
                 ], header_color=ACCENT_PURPLE, body_font=11)

# Network
_box_with_header(slide, Inches(6.7), Inches(4.2), Inches(5.9), Inches(1.8),
                 "Docker Network", [
                     "Network: travel-network (bridge driver)",
                     "All 8 containers on same network",
                     "Internal DNS resolution by service name",
                     "Port mapping: host:container isolation",
                     "No external access to data services",
                 ], header_color=ACCENT_GREEN, body_font=11)

_add_textbox(slide, Inches(0.4), Inches(6.3), Inches(12), Inches(0.4),
             "All data services run as Docker containers with named volumes for persistence and automatic restart policies",
             font_size=11, color=LIGHT_GRAY)


# ====================================================================
# SLIDE 8 – Data Flows
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, DARK_BG)
_slide_title_bar(slide, "Key Data Flows")

flows = [
    ("1. Authentication",
     "User --> Login Form --> POST /api/auth/login --> JWT Tokens --> Zustand authStore",
     ACCENT_BLUE),
    ("2. Flight Search",
     "Search Form --> flightService --> GET /api/flights/search --> FlightSearchTool --> SerpAPI",
     ACCENT_CYAN),
    ("3. AI Trip Planning",
     "AI Planner --> POST /api/agents/plan-trip --> LangGraph Pipeline --> 9 Agents --> Final Recommendation",
     ACCENT_GREEN),
    ("4. Payment",
     "PaymentPage --> paymentService --> POST /api/payments --> Stripe API --> DB --> Email Notification",
     ACCENT_ORANGE),
    ("5. Real-Time Notifications",
     "Event --> Celery --> RabbitMQ --> Django Channels --> Redis --> WebSocket --> UI Toast",
     ACCENT_PURPLE),
    ("6. MCP Agent Communication",
     "Agent A --> Register --> MCP Server --> Redis Pub/Sub --> WebSocket --> Agent B",
     ACCENT_RED),
    ("7. Itinerary Generation",
     "ItineraryPage --> Enhanced Orchestrator --> OpenAI GPT-4o-mini --> PostgreSQL --> PDF Export",
     ACCENT_BLUE),
    ("8. Caching Strategy",
     "Request --> Redis Check --> HIT: return cached  |  MISS: query DB --> store with TTL --> return",
     ACCENT_CYAN),
]

y_start = Inches(1.25)
row_h = Inches(0.72)
for i, (title, desc, clr) in enumerate(flows):
    y = y_start + i * row_h
    # Number/title badge
    _add_rect(slide, Inches(0.4), y, Inches(3.2), Inches(0.58), clr)
    _add_textbox(slide, Inches(0.5), y + Inches(0.1), Inches(3.0), Inches(0.45),
                 title, font_size=12, bold=True, color=WHITE)
    # Flow description
    _add_rect(slide, Inches(3.7), y, Inches(9.2), Inches(0.58), MED_BG,
              border_color=clr, border_width=Pt(1))
    _add_textbox(slide, Inches(3.85), y + Inches(0.1), Inches(8.9), Inches(0.45),
                 desc, font_size=11, color=SOFT_WHITE)

_add_textbox(slide, Inches(0.4), Inches(7.05), Inches(12), Inches(0.3),
             "All flows use JWT authentication  |  Errors handled with DRF exception handler  |  Responses cached where applicable",
             font_size=10, color=LIGHT_GRAY)


# ====================================================================
# SLIDE 9 – External Integrations
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, DARK_BG)
_slide_title_bar(slide, "External Service Integrations")

integrations = [
    ("OpenAI API", [
        "Model: GPT-4o-mini",
        "LLM processing & NLP",
        "Itinerary synthesis",
        "Chat completions",
        "Streaming responses",
    ], ACCENT_PURPLE),
    ("SerpAPI", [
        "Google Flights search",
        "Google Hotels search",
        "Google Local results",
        "Cars & Restaurants",
        "Real-time pricing",
    ], ACCENT_BLUE),
    ("Stripe", [
        "Payment processing",
        "Publishable + Secret keys",
        "Payment intents",
        "Webhook handling",
        "Refund support",
    ], ACCENT_CYAN),
    ("Weather API", [
        "OpenWeatherMap provider",
        "Forecast data",
        "Current conditions",
        "Multi-day outlook",
        "Location-based",
    ], ACCENT_GREEN),
    ("SMTP (Gmail)", [
        "Email notifications",
        "Booking confirmations",
        "Password reset",
        "TLS encryption",
        "Template-based emails",
    ], ACCENT_ORANGE),
    ("ElevenLabs", [
        "Voice TTS output",
        "Audio responses",
        "Multiple voice options",
        "Real-time synthesis",
        "Chat integration",
    ], ACCENT_RED),
]

col_w = Inches(2.0)
col_gap = Inches(0.15)
total_w = len(integrations) * col_w + (len(integrations) - 1) * col_gap
sx = (SW - total_w) // 2

for i, (name, items, clr) in enumerate(integrations):
    x = sx + i * (col_w + col_gap)
    _box_with_header(slide, x, Inches(1.4), col_w, Inches(3.0),
                     name, items, header_color=clr, body_font=10)

# Connection pattern
_add_textbox(slide, Inches(0.4), Inches(4.7), Inches(12.5), Inches(0.4),
             "Integration Pattern:  Backend Service Layer  -->  HTTP Client (requests/httpx)  -->  External API  -->  Response Processing  -->  Cache (Redis)",
             font_size=12, bold=False, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# API key management note
_box_with_header(slide, Inches(0.4), Inches(5.3), Inches(12.5), Inches(1.2),
                 "API Key Management & Security", [
                     "All API keys stored as environment variables (never in code)  |  Docker secrets for production deployment",
                     "Keys: OPENAI_API_KEY, SERPAPI_API_KEY, STRIPE_SECRET_KEY, STRIPE_PUBLISHABLE_KEY, WEATHER_API_KEY, ELEVENLABS_API_KEY",
                     "Rate limiting applied per-service  |  Fallback/retry logic with exponential backoff  |  Usage monitoring & alerts",
                 ], header_color=MED_BG, body_font=10)


# ====================================================================
# SLIDE 10 – Security & DevOps
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, DARK_BG)
_slide_title_bar(slide, "Security & Deployment")

# Security
_box_with_header(slide, Inches(0.4), Inches(1.3), Inches(6.2), Inches(2.8),
                 "Security Configuration", [
                     "SSL Redirect: Enforced in production",
                     "HSTS: max-age 1 year, includeSubDomains",
                     "XSS Filter: Enabled (X-XSS-Protection)",
                     "X-Frame-Options: DENY (clickjacking prevention)",
                     "CSRF: Django middleware + token validation",
                     "Secure Cookies: HttpOnly, Secure, SameSite=Lax",
                     "Content-Type: X-Content-Type-Options: nosniff",
                     "CORS: Whitelist-based origin control",
                     "Rate Limiting: Per-user & per-IP throttling",
                 ], header_color=ACCENT_RED, body_font=10)

# Logging
_box_with_header(slide, Inches(6.9), Inches(1.3), Inches(5.9), Inches(1.3),
                 "Logging & Monitoring", [
                     "RotatingFileHandler: 10 MB x 10 backups",
                     "Console + File output handlers",
                     "Structured logging with levels (DEBUG-CRITICAL)",
                 ], header_color=ACCENT_ORANGE, body_font=10)

# Docker
_box_with_header(slide, Inches(6.9), Inches(2.8), Inches(5.9), Inches(1.3),
                 "Docker Configuration", [
                     "8 containers: nginx, frontend, backend, mcp, celery-worker,",
                     "  celery-beat, postgres, redis, rabbitmq",
                     "Restart policy: unless-stopped  |  Health checks enabled",
                 ], header_color=ACCENT_BLUE, body_font=10)

# Tech stack summary
_box_with_header(slide, Inches(0.4), Inches(4.5), Inches(12.5), Inches(1.6),
                 "Technology Stack Summary", [
                     "Backend: Python 3.11  |  Django 5.0  |  DRF  |  Celery  |  Channels  |  Gunicorn  |  Daphne",
                     "Frontend: Node 20  |  React 18  |  TypeScript 5  |  Vite  |  TailwindCSS 3  |  Zustand  |  React Query",
                     "AI/ML: LangChain  |  LangGraph  |  OpenAI GPT-4o-mini  |  SerpAPI  |  RAG Pipeline",
                     "Infrastructure: Docker Compose  |  Nginx  |  PostgreSQL 15  |  Redis 7  |  RabbitMQ 3.12",
                 ], header_color=ACCENT_CYAN, body_font=11)

# Bottom bar
_add_rect(slide, Inches(0), Inches(6.6), SW, Inches(0.06), ACCENT_BLUE)
_add_textbox(slide, Inches(0.4), Inches(6.75), Inches(12), Inches(0.4),
             "AI Smart Flight Agent  --  Multi-Agent AI Travel Planning System  --  Docker Compose Deployment  --  Production Ready",
             font_size=12, bold=True, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(__file__), "technical-architecture.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
